from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# --- Color Palette (Green Theme) ---
GREEN = RGBColor(0x1B, 0x7A, 0x4C)
GREEN_LIGHT = RGBColor(0x25, 0x96, 0x5C)
GREEN_PALE = RGBColor(0xEC, 0xF8, 0xF1)
NAVY = RGBColor(0x1A, 0x23, 0x32)
NAVY_LIGHT = RGBColor(0x2A, 0x3A, 0x52)
GRAY = RGBColor(0x6B, 0x7A, 0x8D)
GRAY_LIGHT = RGBColor(0x9A, 0xAB, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SKY = RGBColor(0xF0, 0xF5, 0xFA)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)

# --- Setup Presentation ---
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Helper Functions ---
def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width: shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=12, color=NAVY, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def set_cell_format(cell, text, font_size=11, color=NAVY, bold=False, fill_color=None, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    cell.text_frame.word_wrap = True
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
    # Remove default borders for cleaner look
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for border_name in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
        ln = tcPr.find(qn(border_name))
        if ln is None:
            ln = tcPr.makeelement(qn(border_name), {})
            tcPr.append(ln)
        ln.set('w', '0')

# ==========================================
# SLIDE 1: COVER
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5), fill_color=GREEN)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(8.1), Inches(0), Inches(5.233), Inches(7.5), fill_color=NAVY)

add_text(slide, Inches(0.8), Inches(1.5), Inches(6), Inches(0.5), "● PROPOSAL", font_size=12, color=GREEN, bold=True)
add_text(slide, Inches(0.8), Inches(2.2), Inches(6), Inches(1.2), "商品ご提案書", font_size=44, color=NAVY, bold=True)
add_text(slide, Inches(0.8), Inches(3.5), Inches(6), Inches(0.6), "Product Proposal", font_size=18, color=GRAY_LIGHT)

fields = ["案件名", "ご提案先", "提案日", "担当"]
y_pos = 4.8
for field in fields:
    add_text(slide, Inches(0.8), Inches(y_pos), Inches(1.5), Inches(0.4), field, font_size=14, color=GRAY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(y_pos + 0.35), Inches(4.5), Inches(0.02), fill_color=BORDER)
    y_pos += 0.7

add_text(slide, Inches(9.0), Inches(5.5), Inches(3.5), Inches(1), "最適な商品を\nご提案いたします", font_size=14, color=WHITE)

# ==========================================
# SLIDE 2: OVERVIEW
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.4), Inches(0.5), Inches(0.5), fill_color=GREEN)
add_text(slide, Inches(0.65), Inches(0.42), Inches(0.5), Inches(0.5), "1", font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1.3), Inches(0.4), Inches(3), Inches(0.5), "ご提案概要", font_size=24, color=NAVY, bold=True)
add_text(slide, Inches(1.3), Inches(0.85), Inches(3), Inches(0.3), "PROPOSAL OVERVIEW", font_size=10, color=GRAY_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.3), Inches(12), Inches(0.03), fill_color=GREEN)

# Column 1
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.6), Inches(3.8), Inches(5.2), fill_color=SKY, line_color=BORDER, line_width=Pt(1))
add_text(slide, Inches(1.0), Inches(1.9), Inches(3), Inches(0.4), "📄 ご提案の背景", font_size=14, color=NAVY, bold=True)
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.5), Inches(3.2), Inches(1.5), fill_color=WHITE, line_color=BORDER, line_width=Pt(1))
add_text(slide, Inches(1.2), Inches(2.7), Inches(2.8), Inches(1.2), "背景情報をここに記載してください", font_size=12, color=GRAY)

# Column 2
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.6), Inches(3.8), Inches(5.2), fill_color=WHITE, line_color=BORDER, line_width=Pt(1))
add_text(slide, Inches(5.2), Inches(1.9), Inches(3), Inches(0.4), "★ おすすめポイント", font_size=14, color=NAVY, bold=True)

points = ["ポイント１を記載", "ポイント２を記載", "ポイント３を記載"]
py = 2.6
for i, pt in enumerate(points):
    opacity_color = GREEN if i == 0 else GREEN_LIGHT if i == 1 else RGBColor(0x4D, 0xA8, 0x73)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(py), Inches(0.4), Inches(0.4), fill_color=opacity_color)
    add_text(slide, Inches(5.22), Inches(py + 0.05), Inches(0.4), Inches(0.3), str(i + 1), font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, Inches(5.8), Inches(py + 0.05), Inches(2.5), Inches(0.3), pt, font_size=12, color=NAVY)
    py += 0.7

# Column 3
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(1.6), Inches(3.8), Inches(5.2), fill_color=NAVY)
add_text(slide, Inches(9.4), Inches(1.9), Inches(3), Inches(0.4), "🏆 推奨商品", font_size=14, color=WHITE, bold=True)
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.4), Inches(2.6), Inches(3.0), Inches(1.5), fill_color=NAVY_LIGHT)
add_text(slide, Inches(9.6), Inches(2.8), Inches(2.6), Inches(0.3), "RECOMMENDED", font_size=9, color=GRAY_LIGHT, bold=True)
add_text(slide, Inches(9.6), Inches(3.3), Inches(2.6), Inches(0.5), "________________", font_size=16, color=WHITE, bold=True)

add_shape(slide, MSO_SHAPE.OVAL, Inches(9.6), Inches(4.8), Inches(0.5), Inches(0.5), fill_color=GREEN)
add_text(slide, Inches(9.63), Inches(4.83), Inches(0.5), Inches(0.4), "✔", font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(10.3), Inches(4.85), Inches(2), Inches(0.4), "最適な商品をご推奨", font_size=11, color=GRAY_LIGHT)


# ==========================================
# SLIDE 3: PRODUCT DETAIL (Using Table)
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.4), Inches(0.5), Inches(0.5), fill_color=GREEN)
add_text(slide, Inches(0.65), Inches(0.42), Inches(0.5), Inches(0.5), "2", font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1.3), Inches(0.4), Inches(3), Inches(0.5), "商品詳細", font_size=24, color=NAVY, bold=True)
add_text(slide, Inches(1.3), Inches(0.85), Inches(3), Inches(0.3), "PRODUCT DETAILS", font_size=10, color=GRAY_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.3), Inches(12), Inches(0.03), fill_color=GREEN)

# Left: Image Placeholder
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.6), Inches(4.5), Inches(5.2), fill_color=SKY, line_color=BORDER, line_width=Pt(1))
add_text(slide, Inches(1.5), Inches(3.8), Inches(2.5), Inches(1), "[ 商品画像 ]", font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(1.4), Inches(0.4), fill_color=GREEN)
add_text(slide, Inches(0.85), Inches(1.82), Inches(1.4), Inches(0.4), "推奨商品", font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Right: Spec Table
rows, cols = 6, 2
table_shape = slide.shapes.add_table(rows, cols, Inches(5.5), Inches(1.6), Inches(7.2), Inches(3.5))
table = table_shape.table

table.columns[0].width = Inches(2.5)
table.columns[1].width = Inches(4.7)

# Header
set_cell_format(table.cell(0, 0), "📦 商品情報", font_size=14, color=WHITE, bold=True, fill_color=NAVY)
set_cell_format(table.cell(0, 1), "", fill_color=NAVY)

specs = ["商品名", "型番", "参考価格", "納期", "最小ロット"]
for i, spec in enumerate(specs):
    row_idx = i + 1
    bg_color = SKY if i % 2 == 0 else WHITE
    set_cell_format(table.cell(row_idx, 0), spec, font_size=12, color=GRAY, bold=True, fill_color=bg_color)
    set_cell_format(table.cell(row_idx, 1), "________________", font_size=12, color=NAVY, fill_color=bg_color)

# Features Box
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(5.4), Inches(7.2), Inches(1.5), fill_color=GREEN_PALE, line_color=GREEN, line_width=Pt(1))
add_text(slide, Inches(5.9), Inches(5.5), Inches(3), Inches(0.4), "✨ 特徴", font_size=13, color=NAVY, bold=True)
add_text(slide, Inches(5.9), Inches(5.9), Inches(6), Inches(1), "・ 特徴１を記載してください\n・ 特徴２を記載してください\n・ 特徴３を記載してください", font_size=12, color=NAVY)


# ==========================================
# SLIDE 4: COMPARISON (Using Table for structure)
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.4), Inches(0.5), Inches(0.5), fill_color=GREEN)
add_text(slide, Inches(0.65), Inches(0.42), Inches(0.5), Inches(0.5), "3", font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1.3), Inches(0.4), Inches(3), Inches(0.5), "候補商品比較", font_size=24, color=NAVY, bold=True)
add_text(slide, Inches(1.3), Inches(0.85), Inches(3), Inches(0.3), "PRODUCT COMPARISON", font_size=10, color=GRAY_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.3), Inches(12), Inches(0.03), fill_color=GREEN)

cards = [
    {"title": "商品名 ①", "badge": "👑 推奨", "border": GREEN, "bar_width": 3.0, "comment_bg": GREEN_PALE},
    {"title": "商品名 ②", "badge": "候補", "border": BORDER, "bar_width": 2.2, "comment_bg": SKY},
    {"title": "商品名 ③", "badge": "候補", "border": BORDER, "bar_width": 1.6, "comment_bg": SKY}
]

x_pos = 0.6
for card in cards:
    # Card Border
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(1.6), Inches(3.8), Inches(5.2), fill_color=WHITE, line_color=card["border"], line_width=Pt(2))
    # Top Accent Line
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(x_pos), Inches(1.6), Inches(3.8), Inches(0.06), fill_color=card["border"])
    # Image Placeholder
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos + 0.3), Inches(2.0), Inches(3.2), Inches(2.0), fill_color=SKY)
    add_text(slide, Inches(x_pos + 1.2), Inches(2.6), Inches(2), Inches(1), "[ 画像 ]", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)
    # Badge
    badge_color = GREEN if "推奨" in card["badge"] else GRAY
    add_text(slide, Inches(x_pos + 0.3), Inches(4.2), Inches(1.5), Inches(0.3), card["badge"], font_size=11, color=badge_color, bold=True)
    # Title
    add_text(slide, Inches(x_pos + 0.3), Inches(4.6), Inches(3), Inches(0.5), card["title"], font_size=18, color=NAVY, bold=True)
    # Price Bar Background
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos + 0.3), Inches(5.4), Inches(3.2), Inches(0.2), fill_color=BORDER)
    # Price Bar Fill
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos + 0.3), Inches(5.4), Inches(card["bar_width"]), Inches(0.2), fill_color=badge_color)
    # Price Text
    add_text(slide, Inches(x_pos + 0.3), Inches(5.7), Inches(3), Inches(0.3), "¥〇〇〇", font_size=14, color=NAVY, bold=True)
    # Comment Box
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos + 0.3), Inches(6.1), Inches(3.2), Inches(0.6), fill_color=card["comment_bg"])
    add_text(slide, Inches(x_pos + 0.5), Inches(6.15), Inches(2.8), Inches(0.5), "コメントを記載", font_size=11, color=NAVY)

    x_pos += 4.2

# ==========================================
# SLIDE 5: RECOMMENDATION
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.4), Inches(0.5), Inches(0.5), fill_color=GREEN)
add_text(slide, Inches(0.65), Inches(0.42), Inches(0.5), Inches(0.5), "4", font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1.3), Inches(0.4), Inches(3), Inches(0.5), "採用推奨", font_size=24, color=NAVY, bold=True)
add_text(slide, Inches(1.3), Inches(0.85), Inches(3), Inches(0.3), "RECOMMENDATION", font_size=10, color=GRAY_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.3), Inches(12), Inches(0.03), fill_color=GREEN)

# Left Image & Title
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.0), Inches(3.5), Inches(3.5), fill_color=SKY, line_color=GREEN, line_width=Pt(2))
add_text(slide, Inches(2.3), Inches(3.4), Inches(2.5), Inches(1), "[ 採用商品画像 ]", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

add_shape(slide, MSO_SHAPE.OVAL, Inches(4.3), Inches(5.0), Inches(0.7), Inches(0.7), fill_color=GREEN)
add_text(slide, Inches(4.35), Inches(5.05), Inches(0.6), Inches(0.6), "✔", font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(5.9), Inches(3.5), Inches(0.3), "採用推奨商品", font_size=10, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(6.2), Inches(3.5), Inches(0.5), "________________", font_size=20, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

# Right: Reasons Box (Navy)
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.0), Inches(1.6), Inches(6.8), Inches(5.2), fill_color=NAVY)
add_text(slide, Inches(6.5), Inches(2.0), Inches(3), Inches(0.4), "👍 推奨理由", font_size=16, color=WHITE, bold=True)

reasons = [
    ("01", "理由タイトル１", "推奨理由の詳細を記載してください。品質やコスト面での優位性など。"),
    ("02", "理由タイトル２", "推奨理由の詳細を記載してください。納期やサポート面での優位性など。"),
    ("03", "理由タイトル３", "推奨理由の詳細を記載してください。実績や信頼性など。")
]

ry = 2.8
for num, title, desc in reasons:
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(ry), Inches(0.5), Inches(0.5), fill_color=GREEN)
    add_text(slide, Inches(6.52), Inches(ry + 0.05), Inches(0.5), Inches(0.4), num, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, Inches(7.3), Inches(ry), Inches(5), Inches(0.4), title, font_size=14, color=WHITE, bold=True)
    add_text(slide, Inches(7.3), Inches(ry + 0.4), Inches(5), Inches(0.5), desc, font_size=11, color=GRAY_LIGHT)
    if num != "03":
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(ry + 0.95), Inches(5.8), Inches(0.02), fill_color=NAVY_LIGHT)
    ry += 1.3

# ==========================================
# SLIDE 6: SCHEDULE
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.4), Inches(0.5), Inches(0.5), fill_color=GREEN)
add_text(slide, Inches(0.65), Inches(0.42), Inches(0.5), Inches(0.5), "5", font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1.3), Inches(0.4), Inches(3), Inches(0.5), "スケジュール", font_size=24, color=NAVY, bold=True)
add_text(slide, Inches(1.3), Inches(0.85), Inches(3), Inches(0.3), "SCHEDULE", font_size=10, color=GRAY_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.3), Inches(12), Inches(0.03), fill_color=GREEN)

# Timeline Line
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1.79), Inches(2.2), Inches(0.02), Inches(4.5), fill_color=GREEN)

steps = ["📅 提案日", "✅ 採用決定", "🛒 発注", "📦 納品"]
step_descs = ["ご提案書の提出", "商品の採用承認", "正式な発注手続き", "商品の納品完了"]
sy = 2.0
for i, (step, desc) in enumerate(zip(steps, step_descs)):
    circle_color = GREEN if i < 3 else NAVY
    add_shape(slide, MSO_SHAPE.OVAL, Inches(1.5), Inches(sy), Inches(0.6), Inches(0.6), fill_color=circle_color)
    add_text(slide, Inches(1.53), Inches(sy + 0.1), Inches(0.5), Inches(0.4), str(i + 1), font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    add_text(slide, Inches(2.5), Inches(sy), Inches(2), Inches(0.3), step, font_size=16, color=NAVY, bold=True)
    add_text(slide, Inches(4.8), Inches(sy + 0.05), Inches(2), Inches(0.3), "________________", font_size=12, color=GRAY)
    add_text(slide, Inches(2.5), Inches(sy + 0.4), Inches(4), Inches(0.3), desc, font_size=12, color=GRAY)
    sy += 1.2

# Right: Notes
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.6), Inches(5.2), Inches(4.5), fill_color=SKY, line_color=BORDER, line_width=Pt(1))
add_text(slide, Inches(8.0), Inches(2.0), Inches(4), Inches(0.4), "💬 備考", font_size=14, color=NAVY, bold=True)
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(2.6), Inches(4.4), Inches(0.8), fill_color=WHITE, line_color=BORDER, line_width=Pt(1))
add_text(slide, Inches(8.2), Inches(2.7), Inches(4), Inches(0.6), "備考事項を記載してください", font_size=12, color=GRAY)

# Warning Badge
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(6.3), Inches(5.2), Inches(0.8), fill_color=NAVY)
add_text(slide, Inches(7.7), Inches(6.4), Inches(4.8), Inches(0.6), "※ スケジュールは変更される場合がございます。", font_size=11, color=WHITE)

# ==========================================
# SLIDE 7: THANK YOU
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill_color=NAVY)
add_shape(slide, MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.0), Inches(4), Inches(4), fill_color=NAVY_LIGHT)
add_shape(slide, MSO_SHAPE.OVAL, Inches(-1.0), Inches(5.5), Inches(3), Inches(3), fill_color=NAVY_LIGHT)

add_text(slide, Inches(0), Inches(1.5), Inches(13.333), Inches(0.5), "● THANK YOU", font_size=12, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(2.2), Inches(13.333), Inches(1.2), "ありがとうございました", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.0), Inches(3.6), Inches(1.2), Inches(0.04), fill_color=GREEN)
add_text(slide, Inches(0), Inches(3.9), Inches(13.333), Inches(0.8), "ご検討のほど よろしくお願いいたします。", font_size=16, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)

# Company Box
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(5.0), Inches(4.3), Inches(1.8), fill_color=NAVY_LIGHT, line_color=RGBColor(0x3A, 0x4A, 0x62), line_width=Pt(1))
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.2), Inches(5.2), Inches(0.9), Inches(0.6), fill_color=GREEN)
add_text(slide, Inches(6.25), Inches(5.25), Inches(0.9), Inches(0.5), "C", font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(4.5), Inches(5.9), Inches(4.3), Inches(0.3), "Company Name", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(4.5), Inches(6.2), Inches(4.3), Inches(0.3), "Brand Visual", font_size=10, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)

# --- Save File ---
output_path = "Product_Proposal_Green_Theme.pptx"
prs.save(output_path)
print("✅ Xuất file thành công: " + output_path)